"""
多种输入类型的处理器模块。

支持的输入类型：
- 纯文本
- Word文档 (.docx)
- PDF文件 (.pdf)
- Excel文件 (.xlsx)
- PowerPoint文件 (.pptx)
- 图片（使用base64编码用于视觉模型）
- 多文件/文档
"""

import base64
import os
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Optional, Union

try:
    import chardet
except ImportError:
    chardet = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
except ImportError:
    Image = None


class InputType(Enum):
    """支持的输入类型枚举。"""
    TEXT = "text"
    DOCX = "docx"
    PDF = "pdf"
    XLSX = "xlsx"
    PPTX = "pptx"
    IMAGE = "image"
    UNKNOWN = "unknown"


@dataclass
class ProcessedInput:
    """
    处理后输入数据的容器。
    
    属性:
        input_type: 原始输入的类型
        text_content: 提取的文本内容
        images: base64编码的图片列表（用于视觉模型）
        metadata: 关于输入的额外元数据
        source: 原始来源（文件路径或'direct_text'）
    """
    input_type: InputType
    text_content: str
    images: list[dict] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    source: str = "direct_text"
    
    def has_images(self) -> bool:
        """检查输入是否包含图片。"""
        return len(self.images) > 0
    
    def get_combined_content(self) -> str:
        """获取所有文本内容的组合。"""
        return self.text_content


@dataclass
class MultiInput:
    """
    多个处理后输入的容器。
    
    用于用户提供多个文件或混合输入类型时。
    """
    inputs: list[ProcessedInput] = field(default_factory=list)
    
    def add_input(self, processed_input: ProcessedInput) -> None:
        """将处理后的输入添加到集合中。"""
        self.inputs.append(processed_input)
    
    def get_combined_text(self) -> str:
        """获取所有输入的组合文本内容。"""
        texts = []
        for i, inp in enumerate(self.inputs):
            header = f"\n--- 文档 {i + 1} ({inp.source}) ---\n"
            texts.append(header + inp.text_content)
        return "\n".join(texts)
    
    def get_all_images(self) -> list[dict]:
        """获取所有输入中的所有图片。"""
        images = []
        for inp in self.inputs:
            images.extend(inp.images)
        return images
    
    def has_images(self) -> bool:
        """检查是否有任何输入包含图片。"""
        return any(inp.has_images() for inp in self.inputs)


class BaseHandler(ABC):
    """输入处理器的抽象基类。"""
    
    @abstractmethod
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        """检查此处理器是否能处理给定的输入。"""
        pass
    
    @abstractmethod
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        """处理输入并返回处理后的数据。"""
        pass


class TextHandler(BaseHandler):
    """纯文本输入的处理器。"""
    
    SUPPORTED_EXTENSIONS = {'.txt', '.md', '.rst', '.csv', '.json', '.xml', '.yaml', '.yml'}
    
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        if isinstance(input_source, str) and not os.path.exists(input_source):
            # 作为直接文本输入处理
            return True
        
        path = Path(input_source)
        return path.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        if isinstance(input_source, bytes):
            # 检测编码
            if chardet:
                detected = chardet.detect(input_source)
                encoding = detected.get('encoding', 'utf-8')
            else:
                encoding = 'utf-8'
            text = input_source.decode(encoding, errors='replace')
            return ProcessedInput(
                input_type=InputType.TEXT,
                text_content=text,
                source="bytes_input"
            )
        
        if isinstance(input_source, str) and not os.path.exists(input_source):
            # 直接文本输入
            return ProcessedInput(
                input_type=InputType.TEXT,
                text_content=input_source,
                source="direct_text"
            )
        
        # 文件输入
        path = Path(input_source)
        try:
            with open(path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            # 使用检测到的编码重试
            with open(path, 'rb') as f:
                raw = f.read()
            if chardet:
                detected = chardet.detect(raw)
                encoding = detected.get('encoding', 'utf-8')
            else:
                encoding = 'utf-8'
            text = raw.decode(encoding, errors='replace')
        
        return ProcessedInput(
            input_type=InputType.TEXT,
            text_content=text,
            source=str(path),
            metadata={"file_name": path.name, "file_size": path.stat().st_size}
        )


class DocxHandler(BaseHandler):
    """Word文档(.docx)的处理器。"""
    
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        if DocxDocument is None:
            return False
        path = Path(input_source)
        return path.suffix.lower() == '.docx'
    
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        if DocxDocument is None:
            raise ImportError("处理.docx文件需要python-docx库")
        
        path = Path(input_source)
        doc = DocxDocument(path)
        
        # 从段落提取文本
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # 从表格提取文本
        table_texts = []
        for table in doc.tables:
            table_content = []
            for row in table.rows:
                row_content = [cell.text.strip() for cell in row.cells]
                table_content.append(" | ".join(row_content))
            table_texts.append("\n".join(table_content))
        
        text_content = "\n\n".join(paragraphs)
        if table_texts:
            text_content += "\n\n--- 表格 ---\n" + "\n\n".join(table_texts)
        
        return ProcessedInput(
            input_type=InputType.DOCX,
            text_content=text_content,
            source=str(path),
            metadata={
                "file_name": path.name,
                "paragraph_count": len(paragraphs),
                "table_count": len(doc.tables)
            }
        )


class PdfHandler(BaseHandler):
    """PDF文件(.pdf)的处理器。"""
    
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        if PdfReader is None:
            return False
        path = Path(input_source)
        return path.suffix.lower() == '.pdf'
    
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        if PdfReader is None:
            raise ImportError("处理PDF文件需要pypdf库")
        
        path = Path(input_source)
        reader = PdfReader(path)
        
        pages_text = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                pages_text.append(f"--- 第 {i + 1} 页 ---\n{page_text}")
        
        text_content = "\n\n".join(pages_text)
        
        return ProcessedInput(
            input_type=InputType.PDF,
            text_content=text_content,
            source=str(path),
            metadata={
                "file_name": path.name,
                "page_count": len(reader.pages)
            }
        )


class ExcelHandler(BaseHandler):
    """Excel文件(.xlsx)的处理器。"""
    
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        if load_workbook is None:
            return False
        path = Path(input_source)
        return path.suffix.lower() in {'.xlsx', '.xls'}
    
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        if load_workbook is None:
            raise ImportError("处理Excel文件需要openpyxl库")
        
        path = Path(input_source)
        wb = load_workbook(path, data_only=True)
        
        sheets_text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    rows.append(" | ".join(row_values))
            
            if rows:
                sheets_text.append(f"--- 工作表: {sheet_name} ---\n" + "\n".join(rows))
        
        text_content = "\n\n".join(sheets_text)
        
        return ProcessedInput(
            input_type=InputType.XLSX,
            text_content=text_content,
            source=str(path),
            metadata={
                "file_name": path.name,
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames
            }
        )


class PowerPointHandler(BaseHandler):
    """PowerPoint文件(.pptx)的处理器。"""
    
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        if Presentation is None:
            return False
        path = Path(input_source)
        return path.suffix.lower() == '.pptx'
    
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        if Presentation is None:
            raise ImportError("处理PowerPoint文件需要python-pptx库")
        
        path = Path(input_source)
        prs = Presentation(path)
        
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
            
            if slide_content:
                slides_text.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(slide_content))
        
        text_content = "\n\n".join(slides_text)
        
        return ProcessedInput(
            input_type=InputType.PPTX,
            text_content=text_content,
            source=str(path),
            metadata={
                "file_name": path.name,
                "slide_count": len(prs.slides)
            }
        )


class ImageHandler(BaseHandler):
    """
    图片文件的处理器。
    
    将图片转换为base64格式，用于视觉模型。
    """
    
    SUPPORTED_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
    
    def can_handle(self, input_source: Union[str, Path]) -> bool:
        path = Path(input_source)
        return path.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def process(self, input_source: Union[str, Path, bytes]) -> ProcessedInput:
        path = Path(input_source)
        
        # 读取并编码图片
        with open(path, 'rb') as f:
            image_data = f.read()
        
        base64_image = base64.b64encode(image_data).decode('utf-8')
        
        # 确定MIME类型
        extension = path.suffix.lower()
        mime_types = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.webp': 'image/webp',
            '.bmp': 'image/bmp'
        }
        mime_type = mime_types.get(extension, 'image/png')
        
        # 如果PIL可用，获取图片尺寸
        metadata = {"file_name": path.name}
        if Image:
            try:
                with Image.open(path) as img:
                    metadata["width"] = img.width
                    metadata["height"] = img.height
                    metadata["format"] = img.format
            except Exception:
                pass
        
        image_info = {
            "type": "image_url",
            "image_url": {
                "url": f"data:{mime_type};base64,{base64_image}"
            },
            "source": str(path)
        }
        
        return ProcessedInput(
            input_type=InputType.IMAGE,
            text_content=f"[图片: {path.name}]",
            images=[image_info],
            source=str(path),
            metadata=metadata
        )


class InputHandler:
    """
    主输入处理器，根据输入类型委托给特定处理器。
    
    使用方式:
        handler = InputHandler()
        
        # 处理单个文本输入
        result = handler.process_text("一些需求文本")
        
        # 处理单个文件
        result = handler.process_file("requirements.docx")
        
        # 处理多个输入
        result = handler.process_multiple([
            "requirements.docx",
            "screenshots/",
            "额外说明"
        ])
    """
    
    def __init__(self):
        self.handlers: list[BaseHandler] = [
            DocxHandler(),
            PdfHandler(),
            ExcelHandler(),
            PowerPointHandler(),
            ImageHandler(),
            TextHandler(),  # TextHandler应该放在最后作为回退
        ]
    
    def _get_handler(self, input_source: Union[str, Path]) -> BaseHandler:
        """获取输入源的适当处理器。"""
        for handler in self.handlers:
            if handler.can_handle(input_source):
                return handler
        return self.handlers[-1]  # 默认使用TextHandler
    
    def process_text(self, text: str) -> ProcessedInput:
        """
        处理直接文本输入。
        
        参数:
            text: 要处理的文本内容
            
        返回:
            包含文本的ProcessedInput
        """
        return ProcessedInput(
            input_type=InputType.TEXT,
            text_content=text,
            source="direct_text"
        )
    
    def process_file(self, file_path: Union[str, Path]) -> ProcessedInput:
        """
        处理单个文件。
        
        参数:
            file_path: 要处理的文件路径
            
        返回:
            包含提取内容的ProcessedInput
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件未找到: {file_path}")
        
        handler = self._get_handler(path)
        return handler.process(path)
    
    def process_directory(
        self,
        dir_path: Union[str, Path],
        recursive: bool = True
    ) -> MultiInput:
        """
        处理目录中的所有支持文件。
        
        参数:
            dir_path: 目录路径
            recursive: 是否处理子目录
            
        返回:
            包含所有处理文件的MultiInput
        """
        path = Path(dir_path)
        if not path.is_dir():
            raise NotADirectoryError(f"不是目录: {dir_path}")
        
        multi_input = MultiInput()
        
        # 获取所有支持的扩展名
        supported_extensions = set()
        for handler in self.handlers:
            if hasattr(handler, 'SUPPORTED_EXTENSIONS'):
                supported_extensions.update(handler.SUPPORTED_EXTENSIONS)
        supported_extensions.update({'.docx', '.pdf', '.xlsx', '.pptx'})
        
        # 查找并处理文件
        pattern = '**/*' if recursive else '*'
        for file_path in path.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                try:
                    processed = self.process_file(file_path)
                    multi_input.add_input(processed)
                except Exception as e:
                    # 记录错误但继续处理其他文件
                    print(f"警告: 处理 {file_path} 失败: {e}")
        
        return multi_input
    
    def process_multiple(
        self,
        inputs: list[Union[str, Path]]
    ) -> MultiInput:
        """
        处理多个输入（文件、目录或文本）。
        
        参数:
            inputs: 文件路径、目录路径或文本字符串的列表
            
        返回:
            包含所有处理输入的MultiInput
        """
        multi_input = MultiInput()
        
        for input_item in inputs:
            path = Path(input_item) if not isinstance(input_item, Path) else input_item
            
            if path.exists():
                if path.is_file():
                    processed = self.process_file(path)
                    multi_input.add_input(processed)
                elif path.is_dir():
                    dir_result = self.process_directory(path)
                    for inp in dir_result.inputs:
                        multi_input.add_input(inp)
            else:
                # 作为直接文本输入处理
                processed = self.process_text(str(input_item))
                multi_input.add_input(processed)
        
        return multi_input
    
    def process(
        self,
        input_source: Union[str, Path, list[Union[str, Path]]]
    ) -> Union[ProcessedInput, MultiInput]:
        """
        通用处理方法，处理任何输入类型。
        
        参数:
            input_source: 单个输入（文本/文件/目录）或输入列表
            
        返回:
            单个输入返回ProcessedInput，多个输入返回MultiInput
        """
        if isinstance(input_source, list):
            return self.process_multiple(input_source)
        
        path = Path(input_source) if isinstance(input_source, str) else input_source
        
        if isinstance(input_source, str) and not path.exists():
            # 直接文本输入
            return self.process_text(input_source)
        
        if path.is_file():
            return self.process_file(path)
        elif path.is_dir():
            return self.process_directory(path)
        else:
            return self.process_text(str(input_source))
